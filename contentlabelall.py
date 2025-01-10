import os
import requests
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from typing import Optional
from docx import Document
from pdfplumber import open as pdf_open
from pptx import Presentation
from io import BytesIO
import google.generativeai as genai
from dotenv import load_dotenv
import uvicorn
import json
from fastapi.middleware.cors import CORSMiddleware  # Import CORSMiddleware
# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_GEMINI_KEY"))

generation_config = {
    "temperature": 0.3,
    "top_p": 0.95,
    "top_k": 64,
    "max_output_tokens": 8192,
    "response_mime_type": "text/plain",
}

model = genai.GenerativeModel(
    model_name="gemini-1.5-flash",
    generation_config=generation_config,
)

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins, change to specific origins as needed
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods (GET, POST, etc.)
    allow_headers=["*"],  # Allows all headers
)
class FileRequest(BaseModel):
    file_url: str

@app.post("/detect-domain-from-file")
async def detect_domain_from_file(request: FileRequest):
    try:
        # Download the file
        response = requests.get(request.file_url)
        if response.status_code != 200:
            raise HTTPException(status_code=400, detail="Failed to download file.")

        file_content = BytesIO(response.content)
        filename = os.path.basename(request.file_url)
        file_extension = filename.split(".")[-1].lower()

        # Extract content based on file type
        content = ""
        if file_extension == "docx":
            doc = Document(file_content)
            content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == "pdf":
            with pdf_open(file_content) as pdf:
                content = "\n".join(page.extract_text() for page in pdf.pages if page.extract_text())
        elif file_extension == "pptx":
            ppt = Presentation(file_content)
            content = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type.")

        if not content.strip():
            raise HTTPException(status_code=400, detail="Extracted content is empty.")

        # Generate prompt for Gemini
        prompt = f"""
        Analyze the following educational content and determine its subject domain (e.g., Mathematics, Physics, Biology, History, etc.)
        and subdomain (if applicable). Provide a brief explanation for why you classified it as that domain and subdomain.
        Format your response as JSON with three fields: 'domain', 'subdomain', and 'explanation'.
        
        Content: {content}
        """

         # Get response from Gemini
        gemini_response = model.generate_content(prompt)
        
                # Clean and parse the response
        response_text = gemini_response.text.strip()

        # Remove any backticks or extraneous markers
        if response_text.startswith("```json"):
            response_text = response_text[7:]  # Remove the opening ```json
        if response_text.endswith("```"):
            response_text = response_text[:-3]  # Remove the closing ```
        if response_text.startswith("'''"):
            response_text = response_text[3:]  # Remove the opening '''
        if response_text.endswith("'''"):
            response_text = response_text[:-3]  # Remove the closing '''

        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            raise HTTPException(status_code=500, detail="Gemini response is not in valid JSON format.")

        # Return the results with extracted fields
        return {
            "filename": filename,
            "domain": result.get("domain", "Unknown"),
            "subdomain": result.get("subdomain", "Unknown"),
            "explanation": result.get("explanation", "No explanation provided."),
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
